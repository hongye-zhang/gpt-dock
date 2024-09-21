import json
import hashlib
import os
from pptx import Presentation
from pptx.util import Pt
from markdown_it import MarkdownIt
import re
from supabase import create_client, Client

class MarkdownToPPTXConverter:
    # 构造函数，是一个ppt的markdown描述。也可以是ppt的outline，然后再补充细节
    def __init__(self, markdown_text, ppttemplate_path=None):
        self.markdown_text = markdown_text
        self.image_descriptions_prefix = "Image Suggestion: "
        self.slides = []
        if ppttemplate_path is not None:
            if not os.path.exists(ppttemplate_path):
                raise FileNotFoundError(
                    f"Cannot find pptx template file: {ppttemplate_path}"
                )
            self.prs = Presentation(ppttemplate_path)
        else:
            self.prs = Presentation()
        self.md = MarkdownIt()

    # 将细节补充到ppt的outline中。注意细节会替换原outline中的内容
    def merge_content_into_outline(self, supplements: str) -> str:
        """
        Function to insert supplement content into an outline at the corresponding content page,
        preserving the leading whitespaces in the original text for markdown formatting.

        Args:
        - outline (str): The main outline as a string.
        - supplements (str): Supplement content as a string to be inserted.

        Returns:
        - str: The updated outline with supplement content inserted.
        """

        outline = self.markdown_text
        # Split outline and supplements into lines
        outline_lines = outline.split("\n")
        supplements_lines = supplements.split("\n")

        # Create a dictionary to map content pages to their supplement contents
        supplements_dict = {}
        current_page = None
        current_content = []

        # Parse supplements to get the page and content
        for line in supplements_lines:
            if line.startswith("# Content Page"):
                if current_page is not None and current_content:
                    supplements_dict[current_page] = "\n".join(current_content)
                current_page = line  # Keep whitespace intact
                current_content = [line]
            else:
                current_content.append(line)

        if current_page is not None and current_content:
            supplements_dict[current_page] = "\n".join(current_content)

        # Replace the corresponding content in the outline
        updated_outline = []
        for line in outline_lines:
            if line.startswith("# Content Page"):
                page_title = line
                if page_title in supplements_dict:
                    updated_outline.append(supplements_dict[page_title])
                else:
                    updated_outline.append(line)
            else:
                updated_outline.append(line)

        self.markdown_text = "\n".join(updated_outline)

    # 提取image suggestion，变成独立的数据，以方便后续的处理
    def extract_image_suggestions_with_page_numbers(self, ppt_markdown=None):
        if ppt_markdown is None:
            ppt_markdown = self.markdown_text

        # Define a regular expression pattern to capture content page number and image suggestions
        pattern = r"# Content Page (\d+)[\s\S]*?\*\*Image suggestion\*\*:\s*\"(.*?)\""
        # Find all matches in the markdown content
        matches = re.findall(pattern, ppt_markdown)

        # Format the result as a list of tuples with page number and image suggestion
        result = [(f"Content Page {match[0]}", match[1]) for match in matches]
        return result

    # 解析markdown
    import re

    def parse_markdown(self):
        tokens = self.md.parse(self.markdown_text)
        current_slide = None
        i = 0
        list_level = 0

        def start_new_slide(heading_text):
            """根据标题文本确定幻灯片类型并开始新幻灯片"""
            if heading_text == "Title Page":
                slide_type = "title"
            elif heading_text == "Section Page":
                slide_type = "section"
            elif heading_text.startswith("Content Page"):
                slide_type = "content"
            else:
                slide_type = "unknown"
            # 将标题文本作为页名称添加到幻灯片中
            return {"type": slide_type, "page_name": heading_text}

        def add_bullet_point(text, level):
            """将项目符号内容添加到当前幻灯片的内容中"""
            bullet_point = {"text": text, "level": level}
            if "content" not in current_slide:
                current_slide["content"] = []
            current_slide["content"].append(bullet_point)

        def remove_emphasis(text):
            """去除 ** 的标记，但保留其中的内容"""
            return re.sub(r"\*\*(.+?)\*\*", r"\1", text)

        def remove_leading_empty_lines(text):
            """移除文本框中开头的空行"""
            lines = text.splitlines()
            # 过滤掉前导的空行
            non_empty_lines = [line for line in lines if line.strip() != ""]
            return "\n".join(non_empty_lines)

        def clean_text_after_addition(textbox):
            """处理添加到文字框之后的内容，去掉首行空行"""
            # 获取文字框当前的文本内容
            current_text = textbox.text
            # 移除开头的空行
            cleaned_text = remove_leading_empty_lines(current_text)
            # 更新文字框的内容
            textbox.text = cleaned_text

        while i < len(tokens):
            token = tokens[i]

            if token.type == "heading_open":
                # 处理标题，开始新的幻灯片
                content_token = tokens[i + 1]
                if content_token.type == "inline":
                    heading_text = content_token.content.strip()
                    if current_slide:
                        self.slides.append(current_slide)
                    current_slide = start_new_slide(heading_text)
                i += 2

            elif token.type == "bullet_list_open":
                list_level += 1
                i += 1

            elif token.type == "bullet_list_close":
                list_level -= 1
                i += 1

            elif token.type == "list_item_open":
                content_token = tokens[i + 1]
                if content_token.type == "paragraph_open":
                    inline_token = tokens[i + 2]
                    if inline_token.type == "inline":
                        text = inline_token.content.strip()
                        # 去掉 ** 的标记
                        text = remove_emphasis(text)
                        # 匹配 "Key: Value" 形式的内容
                        match = re.match(r"(.+?): (.+)", text)
                        if match:
                            key = match.group(1).strip()
                            value = match.group(2).strip()
                            if key.startswith("Section"):
                                key = "Section"
                            current_slide[key] = value
                        else:
                            add_bullet_point(text, list_level - 1)
                i += 4

            else:
                i += 1

        if current_slide:
            self.slides.append(current_slide)

        # 假设已经将文本内容添加到了 PPT 的文字框中，处理首行空行
        for slide in self.slides:
            if "textbox" in slide:  # 假设文本框对象为 'textbox'
                clean_text_after_addition(slide["textbox"])

    # 创建pptx
    def create_pptx(self, image_descriptions):
        # 定义幻灯片布局
        title_slide_layout = self.prs.slide_layouts.get_by_name(
            "Title Slide"
        )  # 标题幻灯片
        section_slide_layout = self.prs.slide_layouts.get_by_name(
            "Section Header"
        )  # 节标题幻灯片
        content_slide_layout = self.prs.slide_layouts.get_by_name(
            "Title and Content"
        )  # 内容幻灯片
        image_slide_layout = self.prs.slide_layouts.get_by_name(
            "Title with Content and Picture"
        )  # 假设版式5是带图像的内容幻灯片

        for slide_data in self.slides:
            slide_type = slide_data["type"]
            slide_title = slide_data.get("Title", "")
            slide_page_name = slide_data.get("page_name", "")

            # 查找是否有图像描述
            image_description = next(
                (
                    desc
                    for title, desc in image_descriptions
                    if title == slide_page_name
                ),
                None,
            )

            if slide_type == "title":
                slide_layout = title_slide_layout
                slide = self.prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = slide_data.get("Title", "")
                subtitle.text = slide_data.get("Subtitle", "")

            elif slide_type == "section":
                slide_layout = section_slide_layout
                slide = self.prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_data.get("Section", "")

            elif slide_type == "content":
                # 判断是否有图像描述
                if image_description:
                    slide_layout = image_slide_layout  # 使用带图像的版式
                else:
                    slide_layout = content_slide_layout  # 使用标准内容版式

                slide = self.prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_data.get("Title", "")
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame

                # 移除默认段落中的空行，避免首行空白
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    p.text = ""  # 确保第一个段落没有内容

                if "content" in slide_data:
                    first_paragraph = True  # 标记是否是第一个段落
                    for bullet in slide_data["content"]:
                        if first_paragraph:
                            # 处理第一个段落
                            p = tf.paragraphs[0]
                            first_paragraph = False
                        else:
                            # 其他段落正常添加
                            p = tf.add_paragraph()

                        p.text = bullet["text"]
                        p.level = bullet["level"]
                        p.font.size = Pt(18)

                # 如果有图像描述，则将图像描述添加到幻灯片中
                # pptx这个库无法添加alter-text，因此只能添加一个文本框，来说明图片的内容
                if (
                    image_description
                    and len(image_descriptions) > 0
                    and len(slide.shapes.placeholders) > 2
                ):
                    # 由于未知原因，placeholder的编号不是连续的，所以，我们要根据出现的顺序来判定
                    # 不能简单地根据编号
                    for img_placeholder in slide.shapes.placeholders:
                        pidx = img_placeholder.placeholder_format.idx
                        if pidx <= 1:
                            continue
                        # 假设接着就是图片占位符
                        tf_img = img_placeholder.text_frame
                        tf_img.text = self.image_descriptions_prefix + image_description
                        break

            else:
                continue

    def save_pptx(self, filename="output.pptx"):
        self.prs.save(filename)
        with open(filename, 'rb') as file:
            contents = file.read()
        filename = generate_file_hash(contents)
        url: str = 'https://tdklrrxdggwsbfdvtlws.supabase.co'
        key: str = 'eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InRka2xycnhkZ2d3c2JmZHZ0bHdzIiwicm9sZSI6InNlcnZpY2Vfcm9sZSIsImlhdCI6MTcwOTc1MzA3MSwiZXhwIjoyMDI1MzI5MDcxfQ.a8mYI-pyEnmHqj7S30uEpOdIyjKhEbGPu62yTq961eE'
        supabase = create_client(url, key)
        bucket_name: str = "pptx"
        data = supabase.storage.from_(bucket_name).upload('user/' + filename, contents,
                                                          file_options={
                                                              "content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"})
        res = supabase.storage.from_(bucket_name).get_public_url('user/' + filename)
        return res



def markdown_to_json(markdown_text):
    lines = markdown_text.split("\n")
    result = {"title": "", "subtitle": "", "sections": []}
    current_section = None
    current_content = None

    for line in lines:
        if not line or line == "---":
            continue

        if line.startswith("# "):
            result["title"] = line[2:].strip("*")
        elif line.startswith("### ") and not line.startswith("### **Content Page:"):
            result["subtitle"] = line[4:].strip("*")
        elif line.startswith("## **Section Page:"):
            current_section = {"title": line[18:].strip("*"), "pages": []}
            result["sections"].append(current_section)
        elif line.startswith("### **Content Page:"):
            current_content = {
                "title": line[20:].strip("*"),
                "content": [],
                "image_suggestion": "",
                "midjourney_prompt": "",
            }
            current_section["pages"].append(current_content)
        elif line.startswith("- **"):
            current_content["content"].append(
                {"subtitle": line[3:].strip("*"), "points": []}
            )
        elif line.startswith("  - "):
            current_content["content"][-1]["points"].append(line[4:])
        elif line.startswith("**Image Suggestion**:"):
            current_content["image_suggestion"] = line.split(":", 1)[1].strip()
        elif line.startswith("**Midjourney Prompt**:"):
            current_content["midjourney_prompt"] = line.split(":", 1)[1].strip()

    return json.dumps(result, ensure_ascii=False, indent=2)


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title


def add_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    for item in content:
        p = tf.add_paragraph()
        p.text = item["subtitle"]
        p.level = 0
        for point in item["points"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1


def generate_file_hash(file_content):
    """Generate a SHA-256 hash for the given file content."""
    sha256_hash = hashlib.sha256()

    # Assuming the file_content is in bytes
    # If it's not the case, you should convert it to bytes
    sha256_hash.update(file_content)

    return sha256_hash.hexdigest()


def create_ppt_from_json(json_data):
    prs = Presentation()

    data = json.loads(json_data)

    add_title_slide(prs, data["title"], data["subtitle"])

    for section in data["sections"]:
        add_section_slide(prs, section["title"])
        for page in section["pages"]:
            add_content_slide(prs, page["title"], page["content"])

    prs.save("result.pptx")
